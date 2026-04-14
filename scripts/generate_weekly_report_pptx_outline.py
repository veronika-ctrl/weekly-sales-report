#!/usr/bin/env python3
"""
Generate a starter PowerPoint that mirrors the structure of the W13 2026 sales deck:
  - Slide 1: title + date range
  - Slides 2–7: placeholders for app screenshots
  - Slide 8: Markets section title + date range
  - Slides 9–18: more placeholders

Requires: pip install -e ".[slides]"   (python-pptx)

Usage:
  python scripts/generate_weekly_report_pptx_outline.py
  python scripts/generate_weekly_report_pptx_outline.py --out /path/to/out.pptx
"""
from __future__ import annotations

import argparse
from pathlib import Path

try:
    from pptx import Presentation
except ImportError as e:  # pragma: no cover - runtime hint
    raise SystemExit(
        "Missing dependency: install with  pip install -e \".[slides]\"  "
        "or  pip install python-pptx"
    ) from e

ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUT = ROOT / "presentations" / "CDLP_Sales_Report_outline_generated.pptx"

PLACEHOLDER_BODY = (
    "Replace with a screenshot from the weekly report app, or export figures via "
    "scripts/build_slides_markdown.py + saved JSON under reports/."
)


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate weekly report outline .pptx")
    parser.add_argument(
        "--out",
        type=Path,
        default=DEFAULT_OUT,
        help=f"Output path (default: {DEFAULT_OUT})",
    )
    parser.add_argument(
        "--title-line1",
        default="WEEKLY REPORT",
        help="Cover title line 1",
    )
    parser.add_argument(
        "--title-line2",
        default="23 - 29 of March 2026",
        help="Cover subtitle (date range)",
    )
    parser.add_argument(
        "--markets-line1",
        default="Markets",
        help="Markets section title",
    )
    parser.add_argument(
        "--markets-line2",
        default="16 - 22 of March 2026",
        help="Markets section subtitle",
    )
    args = parser.parse_args()

    prs = Presentation()

    # --- Slide 1: title ---
    layout_title = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout_title)
    slide.shapes.title.text = args.title_line1
    slide.placeholders[1].text = args.title_line2

    # --- Slides 2–7: screenshot placeholders ---
    layout_content = prs.slide_layouts[1]
    for n in range(2, 8):
        slide = prs.slides.add_slide(layout_content)
        slide.shapes.title.text = f"Weekly report — slide {n} (screenshot)"
        slide.shapes.placeholders[1].text_frame.text = PLACEHOLDER_BODY

    # --- Slide 8: Markets section ---
    slide = prs.slides.add_slide(layout_title)
    slide.shapes.title.text = args.markets_line1
    slide.placeholders[1].text = args.markets_line2

    # --- Slides 9–18 ---
    for n in range(9, 19):
        slide = prs.slides.add_slide(layout_content)
        slide.shapes.title.text = f"Markets / detail — slide {n} (screenshot)"
        slide.shapes.placeholders[1].text_frame.text = PLACEHOLDER_BODY

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"Wrote {args.out}")


if __name__ == "__main__":
    main()
